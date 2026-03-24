from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors (corrected: NM90 is primary dark bg) ─────────────────────────────
NM90   = RGBColor(0x1D, 0x1D, 0x2C)
NM80   = RGBColor(0x34, 0x32, 0x47)
NM70   = RGBColor(0x55, 0x52, 0x61)
NM60   = RGBColor(0x6B, 0x68, 0x7A)
NM30   = RGBColor(0xB8, 0xB5, 0xC5)
NM10   = RGBColor(0xF0, 0xEC, 0xE5)
NM05   = RGBColor(0xF8, 0xF6, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xFF, 0xB3, 0x2D)
PURPLE = RGBColor(0x87, 0x2D, 0xED)
PURPLE_DARK = RGBColor(0x6A, 0x1F, 0xC2)
GREEN  = RGBColor(0x19, 0xBA, 0x5A)

W, H = 13.33, 7.5
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_color: s.line.color.rgb = line_color; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def circle(slide, l, t, w, h, line_color, lw=1.2):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.background()
    s.line.color.rgb = line_color
    s.line.width = Pt(lw)
    return s

def txt(slide, text, l, t, w, h, font, size, color, align=PP_ALIGN.LEFT, wrap=True, bold=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.bold = bold
    return tb

def eye(slide, text, l, t, w, color):
    return txt(slide, text.upper(), l, t, w, 0.35, "Roboto Mono", 10, color, wrap=False)

def hed(slide, text, l, t, w, h, color, size=44):
    return txt(slide, text, l, t, w, h, "League Gothic", size, color)

def bdy(slide, text, l, t, w, h, color, size=14):
    return txt(slide, text, l, t, w, h, "Roboto", size, color)

def logo(slide, color=WHITE):
    """Bottom-left logo stand-in: ASTRONOMER in Roboto Mono"""
    return txt(slide, "ASTRONOMER", 0.75, 6.95, 3.0, 0.35, "Roboto Mono", 9, color, wrap=False)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (dark, NM90 bg, orbit decoration)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NM90)

# Gold top stripe
box(s, 0, 0, W, 0.08, fill=GOLD)

# Orbit decoration (dark bg only)
circle(s,  8.2, -2.2, 10.0, 10.0, NM80,   0.8)
circle(s,  9.0, -1.3,  7.2,  7.2, PURPLE, 0.4)
circle(s, -2.0,  5.0,  5.0,  5.0, NM80,   0.8)

# Content
box(s, 0.85, 2.6, 0.05, 2.4, fill=GOLD)
eye(s, "Astronomer  .  Platform Overview", 1.1, 2.55, 10, GOLD)
hed(s, "95% of GenAI pilots\nnever reach production.", 1.1, 2.95, 9.5, 2.5, WHITE, 52)
bdy(s, "The data engineering layer is why. Astronomer fixes it.", 1.1, 5.7, 9.0, 0.55, NM10, 15)
logo(s, RGBColor(0x55, 0x52, 0x61))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Big statement (full Purple bg)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, PURPLE)

# Gold top stripe
box(s, 0, 0, W, 0.08, fill=GOLD)

# Darker right column
box(s, 9.8, 0, 3.53, H, fill=PURPLE_DARK)

eye(s, "The problem", 0.85, 0.55, 6, GOLD)
hed(s, "Fragmented pipelines.\nBroken data.\nStalled AI.", 0.85, 0.95, 9.0, 3.5, WHITE, 64)
box(s, 0.85, 4.65, 6.8, 0.05, fill=WHITE)
bdy(s,
    "46% of organizations say Airflow problems halt their entire operation.\n"
    "The teams spending 40% of their time firefighting pipelines are the same ones\n"
    "failing to get AI into production.",
    0.85, 4.85, 8.8, 1.6, RGBColor(0xE0, 0xD0, 0xFF), 14)
logo(s, RGBColor(0xA0, 0x70, 0xFF))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Split: dark left + card grid right (The Solution)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NM10)

# Gold top stripe
box(s, 0, 0, W, 0.08, fill=GOLD)

# Dark left column
box(s, 0, 0.08, 4.5, H - 0.08, fill=NM90)
eye(s, "The solution", 0.5, 0.7, 3.5, GOLD)
hed(s, "UNIFIED\nORCHES\nTRATION.", 0.5, 1.1, 3.8, 4.8, WHITE, 56)
bdy(s, "Built on Apache Airflow.\nThe fastest path to trusted data\nand production AI.", 0.5, 6.05, 3.6, 0.7, NM60, 12)

# 2x2 card grid — tighter, fills the right panel
cards = [
    ("Build",    "AI-native dev tools. Ship DAGs in minutes, not days."),
    ("Run",      "2x parallel load. 30%+ more DAG runs. 3-4x faster task starts."),
    ("Observe",  "Full lineage, alerting, and data quality checks built in."),
    ("Scale",    "From 10 DAGs to 10,000. Multi-cloud. Any team size."),
]
for i, (heading, desc) in enumerate(cards):
    col = i % 2; row = i // 2
    x = 4.75 + col * 4.2; y = 0.15 + row * 3.6
    box(s, x, y, 4.0, 3.35, fill=WHITE)
    box(s, x, y, 0.05, 3.35, fill=PURPLE)
    hed(s, heading, x + 0.2, y + 0.2, 3.6, 0.65, NM80, 26)
    bdy(s, desc, x + 0.2, y + 0.9, 3.6, 2.2, NM70, 14)

logo(s, NM30)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Stats (dark NM90, 3 stat cards)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NM90)

# Gold top stripe
box(s, 0, 0, W, 0.08, fill=GOLD)

eye(s, "Astro vs. the alternatives", 0.8, 0.45, 10, GOLD)
hed(s, "The benchmarks don't lie.", 0.8, 0.8, 9, 1.1, WHITE, 38)

stats = [
    ("2x",   "the concurrency of AWS MWAA\nat equivalent infrastructure"),
    ("30%+", "more DAG runs per dollar\nthan GCP Composer or MWAA"),
    ("85%",  "fewer DAG failures from\nRedis faults vs. Celery"),
]
for i, (num, label) in enumerate(stats):
    x = 0.7 + i * 4.15
    box(s, x, 2.15, 3.85, 4.35, fill=NM80)
    box(s, x, 2.15, 3.85, 0.06, fill=PURPLE)
    hed(s, num, x + 0.25, 2.4, 3.35, 1.7, GOLD, 72)
    bdy(s, label, x + 0.25, 4.2, 3.35, 1.2, NM10, 13)

logo(s, NM60)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Conclusion / CTA (dark, orbit)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s, NM90)

# Gold top stripe
box(s, 0, 0, W, 0.08, fill=GOLD)

# Orbit decoration
circle(s,  7.5, -1.5, 9.0, 9.0, NM80,   0.8)
circle(s,  8.3, -0.7, 6.8, 6.8, GOLD,   0.4)
circle(s, -2.0,  4.5, 5.2, 5.2, PURPLE, 0.5)

box(s, 0.85, 2.0, 0.05, 2.8, fill=GOLD)
eye(s, "Get started", 1.1, 1.95, 9, GOLD)
hed(s, "Stop firefighting.\nStart shipping AI.", 1.1, 2.35, 9.5, 2.2, WHITE, 52)

steps = ["Request a demo", "Start a free trial", "Talk to an engineer"]
for i, step in enumerate(steps):
    x = 1.1 + i * 3.65
    box(s, x, 4.85, 3.3, 0.72, fill=NM80)
    box(s, x, 4.85, 0.05, 0.72, fill=PURPLE)
    bdy(s, step, x + 0.2, 4.97, 3.0, 0.45, WHITE, 14)

bdy(s, "astronomer.io", 1.1, 6.65, 4.0, 0.35, NM60, 11)
logo(s, NM60)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.expanduser("~/Downloads/astro-platform-example-v2.pptx")
prs.save(out)
print(f"Saved: {out}")
